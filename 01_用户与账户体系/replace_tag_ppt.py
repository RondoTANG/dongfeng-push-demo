import shutil
from pptx import Presentation

def replace_text_in_paragraph(paragraph, replacements):
    text = paragraph.text
    changed = False
    for old, new in replacements.items():
        if old in text:
            text = text.replace(old, new)
            changed = True
    if changed:
        if paragraph.runs:
            p_format = paragraph.runs[0].font
            paragraph.text = text
            for run in paragraph.runs:
                run.font.name = p_format.name
                run.font.size = p_format.size
                run.font.bold = p_format.bold
                run.font.italic = p_format.italic
                run.font.color.rgb = p_format.color.rgb if p_format.color.type == 1 else None
        else:
            paragraph.text = text

def generate_ppt(source, target, replacements):
    shutil.copy2(source, target)
    prs = Presentation(target)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                replace_text_in_paragraph(paragraph, replacements)
    prs.save(target)

# Slide 1: 重构思路
replacements_p1 = {
    "P1-3 业务流程方案-标签覆盖池与轮回分发机制": "P1-3 业务流程方案-标签自动化重构与成长数据打通",
    "首创“覆盖池与轮回”派单机制：建立“全员、尖兵队、守卫队”三大覆盖池，引入批次轮回分发机制，解决任务派发过度集中与分配不均问题，保障全体用户的接单机会均等。": "全面打通成长体系数据：系统每日自动计算经验值(XP)，达标后自动为其打上“专家”、“大师”等系统标签。彻底告别人工导表拉群，实现纯自动化的任务分发。",
    "三大基础覆盖池": "当前痛点",
    "(人群建制)": "(依赖人工)",
    "全员池": "每月手动",
    "(大盘流量托底)": "(拉群建单)",
    "尖兵队覆盖池": "效率低下",
    "(月度核心活跃)": "(耗时易错)",
    "守卫队覆盖池": "主观偏差",
    "(高优专属分发)": "(评估不公)",
    "轮回分发调度引擎": "重构方案",
    "(资源均衡调度)": "(自动算力)",
    "自动跨轮回取人机制": "成长数据自动打通",
    "需求人数 > 本轮剩余时，无缝衔接至下轮": "将派单标签与用户的XP段位直接挂钩",
    "防呆互斥与强制保留": "每日自动更新标签",
    "覆盖池严禁人工排除，死守绝对公平红线": "谁达标谁就进池子，无需人工干预",
    "动态名额释放": "去人工化闭环",
    "作业被撤回时，立即释放占用人数回退进度": "全程系统跑数据，对象筛选绝对客观",
    "可视化大盘": "核心价值",
    "(业务监控)": "(业务收益)",
    "轮回进度跟踪": "去中心化",
    "实时显示：": "打破固化圈层",
    "已用人数 / 总人数": "给全员公平机会",
    "批次序列溯源": "精准分发",
    "列表直观显示：": "标签组合策略",
    "当前第 X 次轮回": "提升作业完成率",
    "派单": "分发"
}

source_file = "护卫军标签轮回方案.pptx"
generate_ppt(source_file, "护卫军标签方案_P1_自动化重构.pptx", replacements_p1)

print("PPT updated with plain language!")
