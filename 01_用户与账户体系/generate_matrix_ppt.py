from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# White Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = hex_to_rgb("#F8FAFC")
bg.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
tf = title_box.text_frame
p = tf.add_paragraph()
p.text = "56 项自动标签全景点阵图"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = hex_to_rgb("#0F172A")

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(10), Inches(0.5))
tf = sub_box.text_frame
p = tf.add_paragraph()
p.text = "基于“分发-风控-荣誉”三大标签池的自动化流转"
p.font.size = Pt(14)
p.font.color.rgb = hex_to_rgb("#475569")

pools = [
    {
        "title": "分发标签池 (29)",
        "desc": "一句话描述：精准圈定核心用户，实现高优任务的高效派发。",
        "color": "#2563EB",
        "bg_color": "#EFF6FF",
        "left": Inches(0.5),
        "tags": [
            "大师", "专家", "熟练期", "新手期", "近7天活跃", "近30天高频",
            "微信能力", "抖音能力", "其他平台能力", "能力(原创)", "兴趣偏好",
            "微信视频号", "微博绑定", "知乎绑定", "微信公众号", "头条号绑定",
            "快手绑定", "抖音绑定", "B站绑定", "小红书绑定", "懂车帝绑定",
            "汽车之家绑定", "易车绑定", "尖兵队", "守卫队", "微信绑定", "..."
        ]
    },
    {
        "title": "风险拦截池 (15)",
        "desc": "一句话描述：实时阻断违规劣质用户，保障任务产出的绝对质量。",
        "color": "#DC2626",
        "bg_color": "#FEF2F2",
        "left": Inches(4.766),
        "tags": [
            "数据作假", "敷衍了事", "态度不佳", "高频申诉", "申诉失败高",
            "跨平台申诉", "不看作业要求", "抵触心理", "问题人物", "流失风险用户",
            "近7天未登录", "近30天未登录", "沉默1个月", "沉默2个月", "沉默3个月"
        ]
    },
    {
        "title": "荣誉评优池 (12)",
        "desc": "一句话描述：量化历史贡献与排名，提供精神激励与大盘表彰。",
        "color": "#D97706",
        "bg_color": "#FFFBEB",
        "left": Inches(9.033),
        "tags": [
            "上月积分Top100", "当年命中X次尖兵", "上月命中尖兵", "当年命中X次守卫",
            "上月命中守卫", "当年命中守卫/尖兵", "当年命中X次王牌", "上月命中王牌",
            "近1个月活跃", "近2个月活跃", "近3个月活跃", "120天注册用户"
        ]
    }
]

for pool in pools:
    # Column Background
    col_width = Inches(3.8)
    col = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pool["left"], Inches(1.6), col_width, Inches(5.4))
    col.fill.solid()
    col.fill.fore_color.rgb = hex_to_rgb("#FFFFFF")
    col.line.color.rgb = hex_to_rgb("#E2E8F0")
    col.line.width = Pt(1)

    # Column Title
    t_box = slide.shapes.add_textbox(pool["left"], Inches(1.8), col_width, Inches(0.5))
    tf = t_box.text_frame
    p = tf.add_paragraph()
    p.text = "● " + pool["title"]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(pool["color"])
    p.alignment = PP_ALIGN.CENTER

    # Description
    d_box = slide.shapes.add_textbox(pool["left"] + Inches(0.2), Inches(2.2), col_width - Inches(0.4), Inches(0.6))
    tf = d_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = pool["desc"]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb("#475569")
    
    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pool["left"] + Inches(0.4), Inches(2.8), col_width - Inches(0.8), Pt(1))
    div.fill.solid()
    div.fill.fore_color.rgb = hex_to_rgb("#F1F5F9")
    div.line.fill.background()

    # Draw Tags
    start_y = Inches(3.0)
    start_x = pool["left"] + Inches(0.3)
    
    tag_width = Inches(1.5)
    tag_height = Inches(0.32)
    x_gap = Inches(0.2)
    y_gap = Inches(0.12)
    
    for idx, tag in enumerate(pool["tags"]):
        row = idx // 2
        col_idx = idx % 2
        
        x = start_x + (tag_width + x_gap) * col_idx
        y = start_y + (tag_height + y_gap) * row
        
        # Tag shape
        t_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, tag_width, tag_height)
        t_shape.fill.solid()
        t_shape.fill.fore_color.rgb = hex_to_rgb("#FFFFFF")
        t_shape.line.color.rgb = hex_to_rgb("#CBD5E1")
        
        tf = t_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "● " + tag
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb("#475569")
        p.alignment = PP_ALIGN.CENTER
        
        # Color the dot in the text
        if len(p.runs) > 0:
            p.runs[0].font.color.rgb = hex_to_rgb(pool["color"])

# Watermark
wm_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.4))
p = wm_box.text_frame.add_paragraph()
p.text = "DONGFENG TAG SYSTEM • T+1 00:00 UPDATED"
p.font.size = Pt(9)
p.font.color.rgb = hex_to_rgb("#94A3B8")

prs.save("护卫军标签方案_自动标签点阵图.pptx")
print("Done")
