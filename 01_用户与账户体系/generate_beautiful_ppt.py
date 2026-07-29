from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

prs = Presentation()
# Set 16:9 ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = hex_to_rgb("#111827")
bg.line.fill.background() # No line

# Top Line
top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Cm(0.2))
top_line.fill.solid()
top_line.fill.fore_color.rgb = hex_to_rgb("#E01E2E")
top_line.line.fill.background()

# Title Accent
accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5), Cm(1.5), Cm(0.3), Cm(1.2))
accent.fill.solid()
accent.fill.fore_color.rgb = hex_to_rgb("#E01E2E")
accent.line.fill.background()

# Title Text
title_box = slide.shapes.add_textbox(Cm(2.2), Cm(1.2), Cm(20), Cm(1.5))
tf = title_box.text_frame
p = tf.add_paragraph()
p.text = "标签自动化与分发引擎 (去痛点直达方案)"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = hex_to_rgb("#FFFFFF")

# Subtitle Text
sub_box = slide.shapes.add_textbox(Cm(2.2), Cm(2.5), Cm(20), Cm(1))
tf = sub_box.text_frame
p = tf.add_paragraph()
p.text = "Data-driven tag mapping & automated task distribution engine"
p.font.size = Pt(14)
p.font.color.rgb = hex_to_rgb("#94A3B8")

def create_card(left, top, width, height, title, num, items):
    # Card Background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb("#1A2133")
    card.line.color.rgb = hex_to_rgb("#334155")
    card.line.width = Pt(1)
    
    # Num Box
    num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Cm(1), top + Cm(1), Cm(1.2), Cm(1.2))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = hex_to_rgb("#4C1D95") # fallback purple-ish red
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = hex_to_rgb("#FCA5A5")
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Card Title
    t_box = slide.shapes.add_textbox(left + Cm(2.5), top + Cm(0.8), width - Cm(3), Cm(1.5))
    tf = t_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb("#FFFFFF")
    
    # Line separator
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Cm(1), top + Cm(2.8), width - Cm(2), Pt(1))
    sep.fill.solid()
    sep.fill.fore_color.rgb = hex_to_rgb("#334155")
    sep.line.fill.background()
    
    # Items
    y_offset = top + Cm(3.5)
    for icon, subtitle, desc in items:
        # Subtitle
        s_box = slide.shapes.add_textbox(left + Cm(1), y_offset, width - Cm(2), Cm(1))
        tf = s_box.text_frame
        p = tf.add_paragraph()
        p.text = f"{icon}  {subtitle}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb("#E2E8F0")
        
        # Desc
        d_box = slide.shapes.add_textbox(left + Cm(1.8), y_offset + Cm(0.8), width - Cm(2.8), Cm(2))
        d_box.text_frame.word_wrap = True
        tf = d_box.text_frame
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = hex_to_rgb("#94A3B8")
        
        y_offset += Cm(2.5)

left_items = [
    ("⚡️", "动态算力代替导表拉群", "系统每日自动计算全盘用户的经验值 (XP)，无需运营手工维护名单。"),
    ("🎯", "千人千面精准打标", "一旦XP达标且符合原创/互动偏好，系统即刻为其烙印“专家/大师”标签。"),
    ("🔄", "去人工化实时闭环", "谁达标谁就自动进池子，数据驱动筛选，确保每次任务的受众绝对精准。")
]

right_items = [
    ("🛡️", "三大阶梯覆盖池", "依托标签划定：大师覆盖池（原尖兵队）、专家覆盖池（原守卫队）、全员池。"),
    ("⚖️", "防呆互斥与绝对公平", "严禁人工排除池内名单；达标者自动按批次序列轮回，打破固化垄断。"),
    ("📊", "跨轮补给与名额释放", "需求超额自动跨轮取人；作业被拒/撤回立即释放人数回退进度，全程可视。")
]

# Cards
card_width = Cm(15)
card_height = Cm(12)
top_pos = Cm(4.5)
left_pos_1 = Cm(1.5)
left_pos_2 = Cm(17.5)

create_card(left_pos_1, top_pos, card_width, card_height, "系统级标签自动化 (XP联动)", "01", left_items)
create_card(left_pos_2, top_pos, card_width, card_height, "轮回覆盖池分发引擎", "02", right_items)

# Watermark
wm_box = slide.shapes.add_textbox(Cm(1.5), prs.slide_height - Cm(1.5), Cm(10), Cm(1))
p = wm_box.text_frame.add_paragraph()
p.text = "DONGFENG GROWTH SYSTEM V2.0"
p.font.size = Pt(10)
p.font.color.rgb = hex_to_rgb("#4B5563")

prs.save("护卫军标签方案_高级排版.pptx")
print("Done")
